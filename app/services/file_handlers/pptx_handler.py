"""PowerPoint 文件处理器 (.pptx)

使用 python-pptx 进行文本提取。

注意：仅支持 .pptx 格式（Office 2007+），不支持旧版 .ppt 格式（Office 97-2003）。
"""

from pathlib import Path
from typing import List

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from loguru import logger

from app.config import config
from app.services.file_handlers.base import BaseFileHandler


class PptxHandler(BaseFileHandler):
    """PowerPoint 文件处理器 (.pptx)

    文本提取策略：
    - 逐页遍历幻灯片
    - 每页提取所有形状中的文本（文本框、占位符、表格等）
    - 表格内容用制表符分隔单元格，保留二维结构
    - 每页幻灯片之间用双换行分隔
    """

    @property
    def supported_extensions(self) -> List[str]:
        return ["pptx"]

    def __init__(self):
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.chunk_max_size * 2,
            chunk_overlap=config.chunk_overlap,
            length_function=len,
            is_separator_regex=False,
        )

    def extract_text(self, file_path: str) -> str:
        """使用 python-pptx 从 .pptx 中提取文本"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx 未安装，请运行: pip install python-pptx"
            )

        path = Path(file_path)
        text_parts: List[str] = []

        try:
            prs = Presentation(str(path))

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: List[str] = []

                for shape in slide.shapes:
                    shape_text = self._extract_shape_text(shape)
                    if shape_text:
                        slide_texts.append(shape_text)

                if slide_texts:
                    text_parts.append(f"--- 第 {slide_num} 页 ---\n" + "\n".join(slide_texts))

            logger.info(
                f"PPTX 文本提取完成: {file_path}, "
                f"{len(prs.slides)} 张幻灯片"
            )

        except Exception as e:
            logger.error(f"PPTX 文件打开失败: {file_path}, 错误: {e}")
            raise RuntimeError(f"无法读取 PPTX 文件: {e}") from e

        if not text_parts:
            logger.warning(f"PPTX 文档中未提取到任何文本: {file_path}")
            return ""

        return "\n\n".join(text_parts)

    def split(self, content: str, file_path: str = "") -> List[Document]:
        if not content or not content.strip():
            logger.warning(f"PPTX 文档内容为空: {file_path}")
            return []

        docs = self._splitter.create_documents(
            texts=[content],
            metadatas=[{
                "_source": file_path,
                "_extension": ".pptx",
                "_file_name": Path(file_path).name,
            }],
        )

        logger.info(f"PPTX 分割完成: {file_path} -> {len(docs)} 个分片")
        return docs

    @staticmethod
    def _extract_shape_text(shape) -> str:
        """从形状中提取文本（文本框、占位符、表格、组合）"""
        from pptx.shapes.group import GroupShapes

        # 处理组合形状（递归提取）
        if isinstance(shape, GroupShapes) or hasattr(shape, 'shapes'):
            texts = []
            for child in shape.shapes:
                child_text = PptxHandler._extract_shape_text(child)
                if child_text:
                    texts.append(child_text)
            return "\n".join(texts) if texts else ""

        # 处理表格
        if shape.has_table:
            return PptxHandler._extract_table_text(shape.table)

        # 处理文本框/占位符
        if shape.has_text_frame:
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            return "\n".join(paragraphs) if paragraphs else ""

        return ""

    @staticmethod
    def _extract_table_text(table) -> str:
        """提取表格文本，保留二维结构"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append("\t".join(cells))

        return "\n".join(rows) if rows else ""
